from fastapi import APIRouter, HTTPException, UploadFile, File, Form, Query
from fastapi.responses import JSONResponse
from app.models import QueryRequest
from app.vectorizer import model, get_index_path, encode_batch, encode_text
from pathlib import Path
import docx
import fitz  # PyMuPDF
import os, shutil, pickle, faiss, numpy as np
import google.generativeai as genai
import openpyxl
import pptx

router = APIRouter()

UPLOAD_DIR = "uploads"
TEXT_ROOT = os.path.join(UPLOAD_DIR, "text")
RAW_DIR = os.path.join(UPLOAD_DIR, "raw")
os.makedirs(RAW_DIR, exist_ok=True)

genai.configure(api_key="AIzaSyAEHMRcUjyJovmkRdP7VQsd3wi88ccDBuk")
model_gemini = genai.GenerativeModel("gemini-2.0-flash")
TEXT_DIR = os.path.join("uploads", "text")

def get_text_dir(domain: str):
    path = os.path.join(TEXT_ROOT, domain)
    os.makedirs(path, exist_ok=True)
    return path

# 📥 檔案上傳
@router.post("/upload")
async def upload_file(domain: str = Form(...), file: UploadFile = File(...)):
    try:
        ext = Path(file.filename).suffix.lower()
        raw_path = os.path.join(RAW_DIR, file.filename)
        with open(raw_path, "wb") as f:
            f.write(await file.read())

        content = ""

        if ext == ".txt":
            with open(raw_path, "r", encoding="utf-8") as f:
                content = f.read()

        elif ext == ".pdf":
            doc = fitz.open(raw_path)
            content = "\n".join([page.get_text() for page in doc])

        elif ext == ".docx":
            doc = docx.Document(raw_path)
            content = "\n".join([p.text for p in doc.paragraphs])

        elif ext in [".xlsx", ".xls"]:
            wb = openpyxl.load_workbook(raw_path)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    content += "\t".join([str(cell) if cell else "" for cell in row]) + "\n"

        elif ext in [".pptx", ".ppt"]:
            pres = pptx.Presentation(raw_path)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"

        else:
            raise HTTPException(status_code=400, detail=f"❌ 不支援的檔案格式：{ext}")

        # 儲存為純文字
        text_dir = os.path.join(TEXT_DIR, domain)
        os.makedirs(text_dir, exist_ok=True)
        text_path = os.path.join(text_dir, f"{Path(file.filename).stem}.txt")
        with open(text_path, "w", encoding="utf-8") as f:
            f.write(content.strip())

        return JSONResponse(content={"message": f"✅ 成功處理並儲存為文字：{file.filename}"})

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"上傳失敗：{str(e)}")

# 🧠 向量化（使用 all-MiniLM-L6-v2 本地模型）
@router.post("/vectorize")
def vectorize_texts(domain: str = Query(...)):
    text_dir = get_text_dir(domain)
    texts, metadatas = [], []

    for fname in os.listdir(text_dir):
        if fname.endswith(".txt"):
            path = os.path.join(text_dir, fname)
            with open(path, "r", encoding="utf-8") as f:
                text = f.read().strip()
                if text:
                    texts.append(text)
                    metadatas.append({"filename": fname, "content": text})

    if not texts:
        raise HTTPException(status_code=400, detail="❌ 找不到可向量化的文本")

    embeddings = encode_batch(texts)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(np.array(embeddings).astype("float32"))

    paths = get_index_path(domain)
    faiss.write_index(index, paths["index"])
    with open(paths["meta"], "wb") as f:
        pickle.dump(metadatas, f)

    return {"message": f"✅ [{domain}] 向量化完成，共 {len(texts)} 筆資料"}

# 🤖 查詢（MiniLM 本地模型查詢 + Gemini 回答）
@router.post("/query")
def query_text(request: QueryRequest):
    domain = request.domain.lower()
    question = request.question.strip()
    paths = get_index_path(domain)

    if not os.path.exists(paths["index"]) or not os.path.exists(paths["meta"]):
        raise HTTPException(status_code=400, detail=f"❌ [{domain}] 尚未建立索引")

    try:
        query_vec = encode_text(question).reshape(1, -1).astype("float32")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"❌ 向量轉換失敗：{str(e)}")

    index = faiss.read_index(paths["index"])
    with open(paths["meta"], "rb") as f:
        metadatas = pickle.load(f)

    if query_vec.shape[1] != index.d:
        raise HTTPException(
            status_code=500,
            detail=f"❌ 向量維度不一致：查詢向量 {query_vec.shape[1]}，Index 維度 {index.d}。請重新 /vectorize 建立索引。"
        )

    D, I = index.search(query_vec, k=3)
    matched_contexts = [metadatas[i]["content"] for i in I[0] if i < len(metadatas)]

    prompt = f"以下是相關資料：\n{'\n\n'.join(matched_contexts)}\n\n請根據上面的內容回答問題：{question}，越詳細越好"

    try:
        response = model_gemini.generate_content(prompt)
        return {"answer": response.text, "matched_documents": matched_contexts}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gemini 回覆失敗：{str(e)}")
